from pptx import Presentation
from PIL import Image
import pytesseract
import io


def extract_text_from_pptx(file_path):
    """Extract text from a PPTX file."""
    try:
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from PPTX: {e}")


def extract_text_and_images_from_pptx(file_path):
    """Extract text and perform OCR on images in a PPTX."""
    text = ""
    try:
        presentation = Presentation(file_path)

        for slide_number, slide in enumerate(presentation.slides, start=1):
            text += f"Slide {slide_number}:\n"

            # Extract text from shapes
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"

            # Extract images from the slide
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Shape with embedded image
                    image = shape.image
                    image_stream = io.BytesIO(image.blob)
                    pil_image = Image.open(image_stream)

                    # Perform OCR on the image
                    ocr_text = pytesseract.image_to_string(pil_image)
                    text += f"\nOCR from Image on Slide {slide_number}:\n{ocr_text}\n"
    except Exception as e:
        raise RuntimeError(f"Failed to extract text and images from PPTX: {e}")
    return text
