from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def style_advanced(slide, text, is_front=True):
    """
    An advanced styling function that:
    - Uses a background image or gradient
    - Places a semi-transparent overlay
    - Adds decorative shapes
    - Dynamically adjusts font sizes to prevent overflow
    - Centers text within a shape
    """
    # Slide dimensions (in inches): typically 13.333 x 7.5 for widescreen 16:9
    # slide_width = slide.part.slide_layout.part.shape_tree.part._sld.sldSz.cy
    # Hardcoding known aspect: default slide in python-pptx is 10" x 7.5"
    # We'll assume a 10" x 7.5" slide for simplicity.
    # Adjust if your template differs.
    SLIDE_WIDTH_IN = 10
    SLIDE_HEIGHT_IN = 7.5

    # --- Background Setup ---
    # Option 1: Gradient background
    slide_background = slide.background
    fill = slide_background.fill
    fill.gradient()
    # Two-stop gradient: from a light color to a darker one
    if is_front:
        fill.gradient_stops[0].color.rgb = RGBColor(255, 228, 196)  # Peach
        fill.gradient_stops[1].color.rgb = RGBColor(255, 140, 0)  # Dark Orange
    else:
        fill.gradient_stops[0].color.rgb = RGBColor(173, 216, 230)  # LightBlue
        fill.gradient_stops[1].color.rgb = RGBColor(25, 25, 112)  # MidnightBlue

    # Option 2: If you prefer an image background:
    # slide.shapes.add_picture('background.jpg', 0, 0, width=Inches(SLIDE_WIDTH_IN), height=Inches(SLIDE_HEIGHT_IN))

    # --- Semi-Transparent Overlay ---
    # Add a translucent rectangle overlay to improve text readability
    overlay = slide.shapes.add_picture(
        "semi_transparent_black_overlay.png",
        0,
        0,
        width=Inches(SLIDE_WIDTH_IN),
        height=Inches(SLIDE_HEIGHT_IN),
    )

    # Now you have a uniform black overlay that is 50% transparent.

    # --- Decorative Element ---
    # Add a small decorative shape or icon on a corner
    deco_size = Inches(0.7)
    deco = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(SLIDE_WIDTH_IN - 1),
        Inches(0.5),
        deco_size,
        deco_size,
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White decorative element
    deco.rotation = 45  # Tilt it for style
    deco.line.fill.background()

    # --- Text Content Placement ---
    # We'll use a text box that centers the text.
    text_box_width = Inches(8)
    text_box_height = Inches(4)
    text_box_left = (Inches(SLIDE_WIDTH_IN) - text_box_width) / 2
    text_box_top = (Inches(SLIDE_HEIGHT_IN) - text_box_height) / 2
    text_shape = slide.shapes.add_textbox(
        text_box_left, text_box_top, text_box_width, text_box_height
    )

    # Add text
    text_frame = text_shape.text_frame
    text_frame.word_wrap = True  # Enable line wrapping
    p = text_frame.paragraphs[0]
    p.text = text

    # --- Dynamic Font Sizing ---
    # Heuristic: reduce font size based on text length
    base_font_size = 36
    text_length = len(text)
    # Example heuristic: if text is very long, reduce font size
    if text_length > 200:
        font_size = Pt(24)
    elif text_length > 100:
        font_size = Pt(28)
    else:
        font_size = Pt(base_font_size)

    # Customize text styling
    p.font.name = "Calibri"
    p.font.size = font_size
    p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.font.bold = True

    # Center align the text
    p.alignment = 1  # 0=left, 1=center, 2=right (based on pptx.enum.text import PP_ALIGN if desired)
    text_frame.vertical_anchor = 1  # 0=top, 1=middle, 2=bottom


styles = {
    "advanced": style_advanced,
}


# Main function to create a presentation with a chosen style
def create_pptx_from_flashcards(flashcards, style="classic"):
    # Create a new presentation
    presentation = Presentation()

    # Get the chosen style function
    style_function = styles.get(style, style_advanced)

    for card in flashcards:
        # Add a slide for the 'front' of the card
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[6]
        )  # Blank layout
        style_function(slide, card.get("front", ""), is_front=True)

        # Add a slide for the 'back' of the card
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[6]
        )  # Blank layout
        style_function(slide, card.get("back", ""), is_front=False)

    return presentation
